#!/usr/bin/env python3
"""Chunk and embed TOM supporting documents for the vector index (Brain 2).

Reads .txt / .md / .json / .pdf / .docx / .pptx / .xlsx documents, splits them into overlapping chunks, attaches
metadata (source, locator, classification), and either:
  - writes a chunks JSON artifact (default), or
  - (`--load`) embeds via the Workbench embedding model and upserts into the
    configured vector backend (Azure AI Search prod / local dev).

Per the data-storage strategy, only supporting documents go to the vector store;
the process hierarchy lives in the graph.

Usage:
  python document_to_vector.py --input ./docs_in --output sample_data/finance_vectors.generated.json
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
from pathlib import Path
from typing import Any

_PARA = re.compile(r"\n\s*\n")
DEFAULT_CLASSIFICATION = "KPMG Confidential"

# ---------------------------------------------------------------------------
# Content safety screening (OWASP LLM01 Prompt Injection, LLM04 Data Poisoning,
# LLM08 Vector & Embedding Weaknesses).
#
# Retrieved chunks are later interpolated into synthesis prompts, so ingested
# documents are an indirect prompt-injection surface. Two deterministic
# controls run on every chunk BEFORE it can reach the vector store:
#   1. `sanitize_text` strips zero-width / bidi-control characters that hide
#      instructions from human reviewers (hidden-text attacks).
#   2. `screen_text` flags instruction-like payloads (e.g. "ignore previous
#      instructions") so the knowledge manager must review before loading.
# ---------------------------------------------------------------------------
_HIDDEN_CHARS = re.compile(
    "[\\u200b-\\u200f"  # zero-width space/joiners, LRM/RLM
    "\\u202a-\\u202e"   # bidi embedding/override controls
    "\\u2060-\\u2064"   # word joiner / invisible operators
    "\\ufeff"             # zero-width no-break space (BOM)
    "\\u00ad]"            # soft hyphen
)
_INJECTION_PATTERNS: list[tuple[re.Pattern[str], str]] = [
    (re.compile(r"\bignore\s+(?:all\s+|any\s+)?(?:previous|prior|above|earlier)\s+(?:instructions?|prompts?|rules?)\b", re.I),
     "instruction-override phrase ('ignore previous instructions')"),
    (re.compile(r"\bdisregard\s+(?:all\s+|any\s+)?(?:previous|prior|above|earlier|your)\s+(?:instructions?|prompts?|rules?|guidelines?)\b", re.I),
     "instruction-override phrase ('disregard ... instructions')"),
    (re.compile(r"\bforget\s+(?:all\s+|any\s+)?(?:previous|prior|your)\s+(?:instructions?|rules?|training)\b", re.I),
     "instruction-override phrase ('forget ... instructions')"),
    (re.compile(r"\b(?:reveal|print|repeat|show)\s+(?:your\s+)?system\s+prompt\b", re.I),
     "system-prompt extraction attempt"),
    (re.compile(r"\byou\s+are\s+now\s+(?:a|an|in)\b", re.I),
     "role-reassignment phrase ('you are now ...')"),
    (re.compile(r"\b(?:developer|god|jailbreak)\s*mode\b", re.I),
     "jailbreak keyword"),
    (re.compile(r"\bdo\s+not\s+(?:mention|reveal|tell)\s+(?:this|the\s+user)\b", re.I),
     "concealment instruction addressed to the model"),
    (re.compile(r"\bnew\s+instructions?\s*:", re.I),
     "embedded instruction block ('new instructions:')"),
]


def sanitize_text(text: str) -> str:
    """Remove characters that can hide content from human review."""
    return _HIDDEN_CHARS.sub("", text)


def screen_text(text: str) -> list[str]:
    """Return descriptions of injection-like patterns found in `text`."""
    return [desc for pattern, desc in _INJECTION_PATTERNS if pattern.search(text)]


def screen_chunks(chunks: list[dict[str, Any]]) -> list[str]:
    """Sanitize every chunk in place; tag suspects and return warnings."""
    warnings: list[str] = []
    for chunk in chunks:
        original = str(chunk.get("text", ""))
        cleaned = sanitize_text(original)
        if cleaned != original:
            chunk["text"] = cleaned
            warnings.append(
                f"Hidden/invisible characters stripped from chunk {str(chunk.get('id', '?'))[:12]} "
                f"({chunk.get('source', 'unknown')}, {chunk.get('locator', '?')})"
            )
        findings = screen_text(cleaned)
        if findings:
            chunk["suspect"] = True
            for finding in findings:
                warnings.append(
                    f"Possible prompt-injection payload in chunk {str(chunk.get('id', '?'))[:12]} "
                    f"({chunk.get('source', 'unknown')}, {chunk.get('locator', '?')}): {finding}"
                )
    return warnings



def _content_id(*parts: str) -> str:
    return hashlib.sha256("|".join(parts).encode("utf-8")).hexdigest()


def _metadata(
    classification: str,
    sector: str = "",
    function: str = "",
    technology: str = "",
) -> dict[str, str]:
    return {
        "classification": classification,
        "sector": sector.strip() or "Cross-sector",
        "function": function.strip() or "Unspecified",
        "technology": technology.strip() or "Tech-agnostic",
    }


def _read_segments(path: Path) -> list[dict[str, str]]:
    """Return text segments with locators suitable for citations."""
    suffix = path.suffix.lower()
    if suffix == ".json":
        return []  # handled separately in collect()
    if suffix in {".txt", ".md"}:
        return [{"text": path.read_text(encoding="utf-8", errors="ignore"), "locator": "document"}]
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return [
            {"text": page.extract_text() or "", "locator": f"p. {i + 1}"}
            for i, page in enumerate(reader.pages)
        ]
    if suffix == ".docx":
        from docx import Document

        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                text += "\n" + " | ".join(cell.text.strip() for cell in row.cells)
        return [{"text": text, "locator": "document"}]
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        segments: list[dict[str, str]] = []
        for idx, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            segments.append({"text": "\n".join(texts), "locator": f"slide {idx}"})
        return segments
    if suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        segments = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell) for cell in row]
                if any(v.strip() for v in values):
                    rows.append(" | ".join(values))
            if rows:
                segments.append({"text": "\n".join(rows), "locator": f"sheet {ws.title}"})
        return segments
    if suffix == ".csv":
        with path.open(newline="", encoding="utf-8-sig") as f:
            rows = [" | ".join(row) for row in csv.reader(f)]
        return [{"text": "\n".join(rows), "locator": "csv"}]
    return []


def chunk_text(text: str, target_chars: int = 900, overlap: int = 150) -> list[str]:
    """Paragraph-aware chunking with character overlap to preserve context."""
    paras = [p.strip() for p in _PARA.split(text) if p.strip()]
    chunks: list[str] = []
    buf = ""
    for para in paras:
        if len(buf) + len(para) + 1 <= target_chars:
            buf = f"{buf}\n{para}".strip()
        else:
            if buf:
                chunks.append(buf)
            if len(para) > target_chars:
                # hard-split very long paragraphs
                start = 0
                while start < len(para):
                    chunks.append(para[start : start + target_chars])
                    start += target_chars - overlap
                buf = ""
            else:
                buf = para
    if buf:
        chunks.append(buf)
    # add overlap between adjacent chunks
    if overlap and len(chunks) > 1:
        overlapped = [chunks[0]]
        for i in range(1, len(chunks)):
            tail = chunks[i - 1][-overlap:]
            overlapped.append((tail + " " + chunks[i]).strip())
        chunks = overlapped
    return chunks


def collect(
    input_path: Path,
    classification: str,
    sector: str = "",
    function: str = "",
    technology: str = "",
) -> list[dict[str, Any]]:
    files = [input_path] if input_path.is_file() else sorted(input_path.rglob("*"))
    out: list[dict[str, Any]] = []
    meta = _metadata(classification, sector, function, technology)
    for path in files:
        if path.is_dir():
            continue
        if path.suffix.lower() == ".json":
            data = json.loads(path.read_text(encoding="utf-8"))
            if isinstance(data, list):
                for d in data:
                    text = str(d.get("text", ""))
                    locator = str(d.get("locator", "document"))
                    source = str(d.get("source", path.stem.replace("_", " ").title()))
                    content_hash = _content_id(source, locator, text)
                    d.setdefault("id", content_hash)
                    d.setdefault("content_hash", content_hash)
                    for key, value in meta.items():
                        d.setdefault(key, value)
                    out.append(d)
            continue
        for segment in _read_segments(path):
            text = segment["text"]
            if not text.strip():
                continue
            chunks = chunk_text(text)
            for i, chunk in enumerate(chunks):
                locator = segment["locator"]
                if len(chunks) > 1:
                    locator = f"{locator}, chunk {i + 1}"
                source = path.stem.replace("_", " ").title()
                content_hash = _content_id(
                    str(path.resolve()),
                    locator,
                    chunk,
                    meta["sector"],
                    meta["function"],
                    meta["technology"],
                )
                out.append(
                    {
                        "id": content_hash,
                        "content_hash": content_hash,
                        "text": chunk,
                        "source": source,
                        "locator": locator,
                        **meta,
                    }
                )
    return out


def main() -> None:
    ap = argparse.ArgumentParser(description="Chunk & embed TOM documents for the vector index")
    ap.add_argument("--input", required=True, type=Path, help="File or directory of TOM supporting documents")
    ap.add_argument("--output", type=Path, default=Path("finance_vectors.generated.json"))
    ap.add_argument("--classification", default=DEFAULT_CLASSIFICATION)
    ap.add_argument("--sector", default="", help="Sector metadata applied to every chunk")
    ap.add_argument("--function", default="", help="Function/vertical metadata applied to every chunk")
    ap.add_argument("--technology", default="", help="Technology metadata applied to every chunk")
    ap.add_argument("--load", action="store_true", help="Embed and upsert into the vector backend")
    args = ap.parse_args()

    chunks = collect(args.input, args.classification, args.sector, args.function, args.technology)
    warnings = screen_chunks(chunks)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps(chunks, indent=2))
    print(f"Wrote {args.output} ({len(chunks)} chunks)")
    for w in warnings:
        print(f"WARNING: {w}")

    if args.load:
        import sys

        if warnings:
            sys.exit(
                f"Refusing to load: {len(warnings)} content-safety warning(s). "
                "Review the flagged chunks (marked 'suspect') and re-run, or "
                "load via the ingestion API with fail_on_warnings=false after review."
            )

        sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "backend"))
        from app.clients.vector_store import get_vector_store

        store = get_vector_store()
        n = store.upsert(chunks)
        print(f"Upserted {n} chunks into the vector backend")


if __name__ == "__main__":
    main()
