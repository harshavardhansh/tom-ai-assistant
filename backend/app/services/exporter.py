"""Asset export (POC gotcha #3): a Markdown intermediate is rendered to Word,
PowerPoint, or PDF with consistent structure and a KPMG classification footer.

The process diagram's canonical visual is the SVG returned by the chat API (the
frontend renders it). To keep exports dependency-light and reliable, each format
also embeds the flow as a structured step table/list, which renders identically
everywhere without native SVG rasterization libraries.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Any, Optional

from app.config import get_settings

_BOLD = re.compile(r"\*\*(.+?)\*\*")
_MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}


@dataclass
class ExportResult:
    content: bytes
    filename: str
    media_type: str


def _classification() -> str:
    return get_settings().document_classification


def _parse_blocks(markdown: str) -> list[tuple[str, str]]:
    """Return (kind, text) blocks: kind in {h1,h2,h3,bullet,para}."""
    blocks: list[tuple[str, str]] = []
    for raw in markdown.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("### "):
            blocks.append(("h3", line[4:]))
        elif line.startswith("## "):
            blocks.append(("h2", line[3:]))
        elif line.startswith("# "):
            blocks.append(("h1", line[2:]))
        elif line.lstrip().startswith(("- ", "* ")):
            blocks.append(("bullet", line.lstrip()[2:]))
        else:
            blocks.append(("para", line))
    return blocks


def _flow_rows(flow: dict[str, Any]) -> list[tuple[str, str, str]]:
    by_id = {s["id"]: s for s in flow.get("steps", [])}
    rows = []
    for s in flow.get("steps", []):
        nxt = ", ".join(by_id[n]["name"] for n in s.get("next", []) if n in by_id) or "—"
        rows.append((s.get("name", ""), s.get("role", ""), nxt))
    return rows


class Exporter:
    def export(self, markdown: str, title: str, fmt: str, flow: Optional[dict[str, Any]] = None) -> ExportResult:
        if fmt == "docx":
            return self._docx(markdown, title, flow)
        if fmt == "pptx":
            return self._pptx(markdown, title, flow)
        if fmt == "pdf":
            return self._pdf(markdown, title, flow)
        raise ValueError(f"Unsupported format: {fmt}")

    # -- Word -----------------------------------------------------------
    def _docx(self, markdown: str, title: str, flow: Optional[dict]) -> ExportResult:
        from docx import Document
        from docx.shared import Pt

        doc = Document()
        doc.add_heading(title, level=0)
        for kind, text in _parse_blocks(markdown):
            if kind == "h1":
                doc.add_heading(text, level=1)
            elif kind == "h2":
                doc.add_heading(text, level=2)
            elif kind == "h3":
                doc.add_heading(text, level=3)
            elif kind == "bullet":
                p = doc.add_paragraph(style="List Bullet")
                self._add_runs(p, text)
            else:
                p = doc.add_paragraph()
                self._add_runs(p, text)

        if flow and flow.get("steps"):
            doc.add_heading(f"Process flow: {flow.get('name','')}", level=2)
            table = doc.add_table(rows=1, cols=3)
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            hdr[0].text, hdr[1].text, hdr[2].text = "Step", "Role", "Next"
            for name, role, nxt in _flow_rows(flow):
                cells = table.add_row().cells
                cells[0].text, cells[1].text, cells[2].text = name, role, nxt

        section = doc.sections[0]
        footer = section.footer.paragraphs[0]
        footer.text = _classification()
        for run in footer.runs:
            run.font.size = Pt(8)

        buf = io.BytesIO()
        doc.save(buf)
        return ExportResult(buf.getvalue(), self._fname(title, "docx"), _MIME["docx"])

    def _add_runs(self, paragraph, text: str) -> None:
        pos = 0
        for m in _BOLD.finditer(text):
            if m.start() > pos:
                paragraph.add_run(text[pos : m.start()])
            paragraph.add_run(m.group(1)).bold = True
            pos = m.end()
        if pos < len(text):
            paragraph.add_run(text[pos:])

    # -- PowerPoint -----------------------------------------------------
    def _pptx(self, markdown: str, title: str, flow: Optional[dict]) -> ExportResult:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = _classification()

        body = prs.slides.add_slide(prs.slide_layouts[1])
        body.shapes.title.text = "Response"
        tf = body.placeholders[1].text_frame
        tf.clear()
        first = True
        for kind, text in _parse_blocks(_strip_bold(markdown)):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.text = text
            para.level = {"h1": 0, "h2": 1, "h3": 2, "bullet": 1, "para": 0}.get(kind, 0)

        if flow and flow.get("steps"):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Process flow: {flow.get('name','')}"
            rows = _flow_rows(flow)
            table = slide.shapes.add_table(
                len(rows) + 1, 3, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * (len(rows) + 1))
            ).table
            for j, h in enumerate(("Step", "Role", "Next")):
                table.cell(0, j).text = h
            for i, (name, role, nxt) in enumerate(rows, 1):
                table.cell(i, 0).text = name
                table.cell(i, 1).text = role
                table.cell(i, 2).text = nxt

        buf = io.BytesIO()
        prs.save(buf)
        return ExportResult(buf.getvalue(), self._fname(title, "pptx"), _MIME["pptx"])

    # -- PDF ------------------------------------------------------------
    def _pdf(self, markdown: str, title: str, flow: Optional[dict]) -> ExportResult:
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=18)
        pdf.add_page()

        def line(text: str, h: float) -> None:
            # new_x=LMARGIN keeps the cursor at the left margin so the next
            # full-width multi_cell always has room (avoids fpdf2 zero-width error).
            pdf.multi_cell(0, h, _latin(text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("Helvetica", "B", 16)
        line(title, 9)
        pdf.ln(2)
        for kind, text in _parse_blocks(_strip_bold(markdown)):
            if kind in {"h1", "h2", "h3"}:
                pdf.set_font("Helvetica", "B", 13 if kind == "h1" else 12)
                line(text, 7)
            elif kind == "bullet":
                pdf.set_font("Helvetica", "", 11)
                line(f"  -  {text}", 6)
            else:
                pdf.set_font("Helvetica", "", 11)
                line(text, 6)
            pdf.ln(1)

        if flow and flow.get("steps"):
            pdf.set_font("Helvetica", "B", 12)
            line(f"Process flow: {flow.get('name','')}", 7)
            pdf.set_font("Helvetica", "", 10)
            for name, role, nxt in _flow_rows(flow):
                line(f"  -  {name} | {role} -> {nxt}", 6)

        pdf.set_y(-15)
        pdf.set_font("Helvetica", "I", 8)
        pdf.cell(0, 8, _latin(_classification()), align="C")

        out = pdf.output()
        return ExportResult(bytes(out), self._fname(title, "pdf"), _MIME["pdf"])

    @staticmethod
    def _fname(title: str, ext: str) -> str:
        slug = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-") or "tom-response"
        return f"{slug}.{ext}"


def _strip_bold(markdown: str) -> str:
    return _BOLD.sub(r"\1", markdown)


def _latin(text: str) -> str:
    """fpdf2 core fonts are latin-1; replace unsupported chars gracefully."""
    return text.encode("latin-1", "replace").decode("latin-1")
